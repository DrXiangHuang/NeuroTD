# neurotd_xy_curves_freeform_v2.py
# Generates a PPTX with editable x (red) and y (blue) waveforms as freeform polylines.
# Requires: python-pptx >= 0.6.22 (1.0.2 is fine)
"""Colors: x is (255,0,0) red; y is (49,141,231) blue.
Line width: 1.75 pt (half of 3.5 pt).
Plot width: 4.3″ (half of 8.6″).
"""

import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Try optional enums (not strictly needed in this version)
try:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
except Exception:
    MSO_SHAPE = None
    MSO_AUTO_SHAPE_TYPE = None


# -----------------------------
# 1) Signal & delay construction
# -----------------------------
def make_signals_with_edited_delays(n=1200, seed=1):
    """
    Build x(t) and y(t) where y is a time-warped/shifted version of x.
    We compute the original delays at three spike centers, then modify:
      d1' = d1 / 4
      d2' = -d2
      d3' = d3 * 0.8
    A smooth piecewise-linear delay field d_edit(t) is used to generate y.
    .. math::
        x(t) = 0.1 \sin(2\pi f t) \;+\;
               \sum_{i=1}^3 0.35 \exp\!\left(
                   -\frac{(t-c_i)^2}{w_i^2}
               \right),

    where f=4 Hz, centers c_i and widths w_i define three spike-like events.

    A time-varying delay curve \Delta(t) is sampled at these centers:

    .. math::

        d_i = \Delta(c_i), \quad i=1,2,3.

    The edited delays are

    .. math::

        d_1' = \tfrac{1}{4} d_1, \quad
        d_2' = -d_2, \quad
        d_3' = 0.8\, d_3.

    A smooth interpolant \tilde{\Delta}(t) is constructed that passes
    through (c_i, d_i') and boundary endpoints. The delayed signal is then

    .. math::

        y(t) = 0.1 \sin\!\big(2\pi f (t - \tilde{\Delta}(t))\big)
               \;+\;
               \sum_{i=1}^3 0.35 \exp\!\left(
                   -\frac{((t-\tilde{\Delta}(t))-c_i)^2}{w_i^2}
               \right).

    Both x(t) and y(t) are normalized to [-1,1].
    """
    rng = np.random.default_rng(seed)
    t = np.linspace(0.0, 1.0, n)  # normalized "time" axis

    # x: base sinusoid + three spikes (centers fixed to define "three delays")
    spike_centers = np.array([0.22, 0.56, 0.83])
    spike_widths = np.array([0.018, 0.022, 0.014])
    x = 0.10 * np.sin(2 * np.pi * 4.0 * t)
    for c, w in zip(spike_centers, spike_widths):
        x += 0.35 * np.exp(-(((t - c) / w) ** 2))

    # ORIGINAL delay field (in "seconds" of normalized t): global + slow variation
    base_offset = 0.06
    var_shift = 0.035 * np.sin(2 * np.pi * 1.2 * t)  # smooth, time-varying
    d_orig = base_offset + var_shift  # effective delay(t)

    # Sample the three "current delays" at spike centers
    d1, d2, d3 = np.interp(spike_centers, t, d_orig)

    # Compute edited delays
    d1_edit = d1 / 4.0
    d2_edit = -d2
    d3_edit = d3 * 0.8

    # Build a smooth target delay curve passing through the edited values
    # Add endpoints to keep the curve stable at boundaries
    knots_x = np.array([0.0, spike_centers[0], spike_centers[1], spike_centers[2], 1.0])
    knots_y = np.array([d1_edit, d1_edit, d2_edit, d3_edit, d3_edit])
    d_edit = np.interp(t, knots_x, knots_y)

    # Build y by applying the edited delay (time shift) to x's components
    # For clarity we reconstruct y with the same recipe as x but time-shifted by d_edit
    t_y = (t - d_edit) % 1.0
    y = 0.10 * np.sin(2 * np.pi * 4.0 * t_y)
    for c, w in zip(spike_centers, spike_widths):
        y += 0.35 * np.exp(-(((t_y - c) / w) ** 2))

    # Normalize both to [-1, 1] for neat plotting
    x /= np.max(np.abs(x)) + 1e-9
    y /= np.max(np.abs(y)) + 1e-9

    # Report the delays (original and edited) at the three spike centers
    d_orig_centers = np.array([d1, d2, d3])
    d_edit_centers = np.array([d1_edit, d2_edit, d3_edit])

    return t, x, y, spike_centers, d_orig_centers, d_edit_centers, d_edit


# -----------------------------
# 2) PowerPoint helpers
# -----------------------------
def add_text(slide, left, top, width, height, text, font_size_pt=20, bold=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    return tb


def add_axes_box(slide, left, top, width, height, line_width_pt=1.25):
    # Simple rectangle to suggest an axes box
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height) if MSO_SHAPE else None
    if rect:
        rect.fill.background()
        rect.line.width = Pt(line_width_pt)
    return rect


def build_freeform(shapes, start_x, start_y):
    """
    Handle both possible signatures of build_freeform:
      - build_freeform(start_x, start_y)
      - build_freeform(shape_type, start_x, start_y)
    """
    try:
        return shapes.build_freeform(int(start_x), int(start_y))
    except TypeError:
        # Older signature fallback (some builds)
        return shapes.build_freeform(MSO_AUTO_SHAPE_TYPE.RECTANGLE, int(start_x), int(start_y))


def add_polyline_freeform(slide, xnorm, ynorm, left, top, width, height, rgb=(255, 0, 0), line_width_pt=1.75):
    """
    Draw a *single* editable polyline using FreeformBuilder.add_line_segments.
    xnorm, ynorm in [0,1]; y=0 at bottom.
    """
    X_abs = left + (xnorm * width)
    Y_abs = top + ((1.0 - ynorm) * height)
    pts_local = [(int(X_abs[i] - left), int(Y_abs[i] - top)) for i in range(len(X_abs))]
    ff = build_freeform(slide.shapes, pts_local[0][0], pts_local[0][1])
    ff.add_line_segments(pts_local[1:], close=False)
    shape = ff.convert_to_shape(int(left), int(top))
    shape.fill.background()
    shape.line.color.rgb = RGBColor(*rgb)
    shape.line.width = Pt(line_width_pt)
    return shape


def add_xy_slide(prs, title, t, x, y, show_x=True, show_y=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6), title, font_size_pt=28, bold=True)

    # >>> Change 3: make time axis half width (8.6" -> 4.3")
    left, top, width, height = Inches(0.8), Inches(1.3), Inches(4.3), Inches(4.4)

    add_axes_box(slide, left, top, width, height, line_width_pt=1.0)
    add_text(slide, left, top + height + Inches(0.05), width, Inches(0.4), "time", font_size_pt=16)
    add_text(slide, left - Inches(0.55), top + Inches(1.7), Inches(0.5), Inches(0.6), "amplitude", font_size_pt=14)

    # Normalize X coords to [0,1]
    xnorm = (t - t.min()) / (t.max() - t.min())

    # Map y from [-1,1] into [0,1] and keep vertical margins
    def map_y(sig):
        yn = (sig + 1.0) / 2.0
        return 0.15 + 0.70 * yn

    # >>> Change 1 & 2: x is RED, y is BLUE; line width = half (1.75 pt)
    if show_x:
        add_polyline_freeform(
            slide, xnorm, map_y(x), left, top, width, height, rgb=(255, 0, 0), line_width_pt=1.75
        )  # red x
    if show_y:
        add_polyline_freeform(
            slide, xnorm, map_y(y), left, top, width, height, rgb=(49, 141, 231), line_width_pt=1.75
        )  # blue y

    return slide


def add_xy_delay_slide(prs, title, t, x, y, d_edit):
    """Fourth slide: plot x (red), y (blue), and d_edit (black) together."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6), title, font_size_pt=28, bold=True)

    left, top, width, height = Inches(0.8), Inches(1.3), Inches(4.3), Inches(4.4)
    add_axes_box(slide, left, top, width, height, line_width_pt=1.0)
    add_text(slide, left, top + height + Inches(0.05), width, Inches(0.4), "time", font_size_pt=16)
    add_text(
        slide,
        left - Inches(0.55),
        top + Inches(1.7),
        Inches(0.5),
        Inches(0.6),
        "amplitude / (scaled delay)",
        font_size_pt=14,
    )

    xnorm = (t - t.min()) / (t.max() - t.min())

    def map_y(sig):
        yn = (sig + 1.0) / 2.0
        return 0.15 + 0.70 * yn

    # Normalize delay to [-1,1] so it overlays nicely; handle constant edge-case
    dmin, dmax = float(np.min(d_edit)), float(np.max(d_edit))
    if dmax - dmin < 1e-12:
        d_norm = np.zeros_like(d_edit)
    else:
        d_norm = 2.0 * (d_edit - (dmin + dmax) / 2.0) / (dmax - dmin)  # [-1,1]

    # Draw x (red), y (blue), and d_edit (black)
    add_polyline_freeform(slide, xnorm, map_y(x), left, top, width, height, rgb=(255, 0, 0), line_width_pt=1.75)
    add_polyline_freeform(slide, xnorm, map_y(y), left, top, width, height, rgb=(49, 141, 231), line_width_pt=1.75)
    add_polyline_freeform(slide, xnorm, map_y(d_norm), left, top, width, height, rgb=(0, 0, 0), line_width_pt=1.75)
    return slide


# -----------------------------
# 3) Run & save
# -----------------------------
if __name__ == "__main__":
    t, x, y, spike_centers, d_orig_centers, d_edit_centers, d_edit = make_signals_with_edited_delays()

    # Report the three delays before/after edits
    # Here "time" is normalized to [0,1]; think of it as seconds for display.
    print("Spike centers (t):", spike_centers)
    print("Original delays at centers: ", [f"{v:.6f} s ({v*1000:.3f} ms)" for v in d_orig_centers])
    print("Edited delays at centers:   ", [f"{v:.6f} s ({v*1000:.3f} ms)" for v in d_edit_centers])

    prs = Presentation()
    add_xy_slide(prs, "Signals x and y (together)", t, x, y, True, True)
    add_xy_slide(prs, "x only (red)", t, x, y, True, False)
    add_xy_slide(prs, "y only (blue)", t, x, y, False, True)
    add_xy_delay_slide(prs, "x (red), y (blue), and edited delay d_edit (black)", t, x, y, d_edit)

    out_path = "neurotd_xy_curves_freeform_v3.pptx"
    prs.save(out_path)
    print(f"Saved PowerPoint to: {out_path}")
